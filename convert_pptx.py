"""Convert Ch3 PPTX files to markdown."""
import os
from pptx import Presentation
from pptx.util import Pt

PPTX_FILES = [
    "Material/Ch3-01-02-overviewGraphCoverag (1).pptx",
    "Material/Ch3-03-sourceCode(1).pptx",
    "Material/Ch3-04-design(1).pptx",
    "Material/Ch3-05-06-spec&amp.pptx",
]

OUTPUT_DIR = "Material"


def extract_text_from_shape(shape):
    """Extract text from a shape, preserving structure."""
    lines = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append("")
                continue

            # Detect heading level by font size
            font_size = None
            bold = False
            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size
                if run.font.bold:
                    bold = True

            level = 0
            if font_size:
                pt = font_size / 12700  # EMU to PT
                if pt >= 36:
                    level = 1
                elif pt >= 28:
                    level = 2
                elif pt >= 20:
                    level = 3
                elif pt >= 16:
                    level = 4

            if level > 0:
                lines.append(f"{'#' * level} {text}")
            elif bold and len(text) < 100:
                lines.append(f"**{text}**")
            else:
                lines.append(text)

    elif shape.has_table:
        table = shape.table
        # Header row
        header = []
        for cell in table.rows[0].cells:
            header.append(cell.text.strip())
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for idx, row in enumerate(table.rows):
            if idx == 0:
                continue
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            lines.append("| " + " | ".join(row_data) + " |")

    return lines


def convert_pptx_to_md(pptx_path, output_path):
    """Convert a PPTX file to markdown."""
    prs = Presentation(pptx_path)
    all_lines = []

    for i, slide in enumerate(prs.slides):
        slide_lines = []
        for shape in slide.shapes:
            shape_lines = extract_text_from_shape(shape)
            slide_lines.extend(shape_lines)

        # Filter empty slides
        content = "\n".join(slide_lines).strip()
        if not content:
            continue

        all_lines.append(f"<!-- Slide {i+1} -->")
        all_lines.append("")
        all_lines.extend(slide_lines)
        all_lines.append("")
        all_lines.append("---")
        all_lines.append("")

    md_content = "\n".join(all_lines)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md_content)

    print(f"Converted: {pptx_path} -> {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    for pptx_file in PPTX_FILES:
        if not os.path.exists(pptx_file):
            print(f"SKIP (not found): {pptx_file}")
            continue
        base = os.path.splitext(os.path.basename(pptx_file))[0]
        # Clean up filename
        base = base.replace(" (1)", "").replace("&amp", "").replace("&", "")
        output_path = os.path.join(OUTPUT_DIR, f"{base}.md")
        convert_pptx_to_md(pptx_file, output_path)

    print("\nAll done!")
